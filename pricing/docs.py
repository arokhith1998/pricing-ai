"""Document ingest for Ask-your-Pricekeel.

Parses PDF / DOCX / MD / TXT into paragraph-aware chunks and embeds them
locally with fastembed (the same model used by the column mapper). Chunks are
returned with their source page + filename so answers can cite them precisely.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from io import BytesIO
from typing import Optional

import numpy as np


# ---- dataclasses -----------------------------------------------------------

@dataclass
class Chunk:
    doc_id: str
    doc_name: str
    chunk_idx: int
    text: str
    page: Optional[int] = None
    # Set after embedding; kept off the wire in API responses.
    embedding: Optional[np.ndarray] = field(default=None, repr=False)


# ---- parsers ---------------------------------------------------------------

def parse_pdf(data: bytes) -> list[tuple[int, str]]:
    """[(page_num, page_text)] for each page."""
    from pypdf import PdfReader
    reader = PdfReader(BytesIO(data))
    return [(i + 1, (p.extract_text() or "").strip()) for i, p in enumerate(reader.pages)]


def parse_docx(data: bytes) -> list[tuple[int, str]]:
    """DOCX has no reliable page concept — return one big block as page 1."""
    from docx import Document as DocxDocument
    d = DocxDocument(BytesIO(data))
    text = "\n".join(p.text for p in d.paragraphs if p.text.strip())
    return [(1, text)]


def parse_text(data: bytes) -> list[tuple[int, str]]:
    return [(1, data.decode("utf-8", errors="ignore"))]


def parse_xlsx(data: bytes) -> list[tuple[int, str]]:
    """One block per worksheet — useful for ratecards, discount matrices,
    commission plans, etc. Rows are joined with ' | ' so the LLM can see the
    columnar shape."""
    import openpyxl
    wb = openpyxl.load_workbook(BytesIO(data), data_only=True)
    out: list[tuple[int, str]] = []
    for i, name in enumerate(wb.sheetnames, start=1):
        ws = wb[name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        if rows:
            out.append((i, f"Sheet: {name}\n" + "\n".join(rows)))
    return out


def parse_pptx(data: bytes) -> list[tuple[int, str]]:
    """One block per slide — useful for sales decks and pricing strategy
    presentations. Pulls text from every shape that has any."""
    from pptx import Presentation
    prs = Presentation(BytesIO(data))
    out: list[tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            txt = getattr(shape, "text", "")
            if isinstance(txt, str) and txt.strip():
                texts.append(txt)
        if texts:
            out.append((i, "\n".join(texts)))
    return out


def parse_and_chunk(name: str, data: bytes,
                    max_chars: int = 1200, overlap: int = 200) -> list[tuple[int, str]]:
    """Dispatch by extension and chunk. Returns [(page, chunk_text), ...]."""
    lname = name.lower()
    if lname.endswith(".pdf"):
        pages = parse_pdf(data)
    elif lname.endswith(".docx"):
        pages = parse_docx(data)
    elif lname.endswith(".xlsx"):
        pages = parse_xlsx(data)
    elif lname.endswith(".pptx"):
        pages = parse_pptx(data)
    elif lname.endswith((".md", ".txt")):
        pages = parse_text(data)
    else:
        raise ValueError(f"Unsupported file type: {name}")
    out: list[tuple[int, str]] = []
    for page_num, page_text in pages:
        for c in _chunk_text(page_text, max_chars=max_chars, overlap=overlap):
            out.append((page_num, c))
    return out


def _chunk_text(text: str, max_chars: int = 1200, overlap: int = 200) -> list[str]:
    """Greedy windowing on whitespace-normalized text, back-stepping to a
    sentence boundary when one is available within the back half of the window.
    """
    text = re.sub(r"\s+", " ", text).strip()
    if len(text) <= max_chars:
        return [text] if text else []
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + max_chars, len(text))
        if end < len(text):
            sb = text.rfind(". ", start, end)
            if sb > start + max_chars // 2:
                end = sb + 1
        chunks.append(text[start:end].strip())
        if end == len(text):
            break
        start = max(end - overlap, start + 1)
    return [c for c in chunks if c]


# ---- embeddings (fastembed; shared model with the column mapper) -----------

_EMBED_MODEL = None


def _model():
    global _EMBED_MODEL
    if _EMBED_MODEL is None:
        from fastembed import TextEmbedding
        _EMBED_MODEL = TextEmbedding(model_name="BAAI/bge-small-en-v1.5")
    return _EMBED_MODEL


def embed_texts(texts: list[str]) -> np.ndarray:
    """(N, D) L2-normalized float32."""
    m = _model()
    embs = np.array(list(m.embed(texts)), dtype=np.float32)
    norms = np.linalg.norm(embs, axis=1, keepdims=True)
    norms = np.where(norms == 0, 1.0, norms)
    return embs / norms


def embed_query(query: str) -> np.ndarray:
    """(D,) L2-normalized."""
    m = _model()
    q = np.array(next(iter(m.embed([query]))), dtype=np.float32)
    n = np.linalg.norm(q)
    return q / (n if n else 1.0)
